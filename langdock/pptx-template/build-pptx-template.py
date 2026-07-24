#!/usr/bin/env python3
# ============================================================
# NATIVER PPTX-MASTER — aus vermessener Deck-Geometrie
#
# Baut aus geometry.json (extract-geometry.js), den text-freien
# Hintergrund-Renderings (Moment-Slides) und den 2x-Screenshots
# (assets/shots/, für Grafik-Ausschnitte) eine echte PowerPoint-
# Datei: Flächen als abgerundete Rechtecke, SVG-/Gradient-
# Grafiken als Bildausschnitte, Textfelder mit Zeilenabstand und
# Laufweite, Hairline-Shapes, Wortmarken — und eingebetteten
# Schriften (Inter, Inter Light, Inter SemiBold, JetBrains Mono
# als TTF-Instanzen aus assets/fonts): PowerPoint substituiert
# sonst auf Rechnern ohne Inter alles durch Calibri (Befund
# Nils, 23.07.2026).
#
# Koordinaten: 1440x810 px = 13,333x7,5 Zoll -> 108 px/Zoll.
# Schriftgrad: pt = px * 2/3. Laufweite: a:spc in 1/100 pt.
#
# Nutzung:
#   node langdock/pptx-template/extract-geometry.js
#   python3 langdock/pptx-template/build-pptx-template.py
#   (braucht: python-pptx, Pillow, fontTools, brotli)
# ============================================================

import json, os, re, zipfile
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

def flat(shp):
    """Theme-Style vom Autoshape strippen: der Style-Verweis bringt
    den Theme-Schatten mit (sichtbar u. a. im LibreOffice-Render) —
    Flächen und Hairlines laufen flach."""
    st = shp._element.find(qn('p:style'))
    if st is not None:
        shp._element.remove(st)
    shp.fill.solid()
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, 'bridgemaker-slides-template.pptx')
EMBED = os.path.join(HERE, 'assets/embed-fonts')
SHOTS = os.path.join(HERE, 'assets/shots')
PX = 914400 / 108                    # EMU pro Pixel

def emu(v): return Emu(int(v * PX))

def rgb(css):
    m = re.match(r'rgba?\((\d+),\s*(\d+),\s*(\d+)', css or '')
    return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3))) if m else RGBColor(0x1C, 0x1C, 0x1E)

def crop_alpha(src, dst):
    img = Image.open(src).convert('RGBA')
    box = img.getbbox()
    img.crop(box).save(dst)
    return dst

# ---------- Schrift-Instanzen für die Einbettung ----------
# PowerPoint bettet nur TTF ein; die vendorten Webfonts sind EIN
# variabler Font — hier werden statische Schnitte instanziert.
# Gewichts-Mapping der Textläufe: <=300 Inter Light, 500-699
# Inter SemiBold (eigene Familien, PP kennt nur Regular/Bold),
# >=700 Inter Bold, sonst Inter Regular.
EMBED_FILES = [
    ('Inter',          'Inter-Regular.ttf',        'Inter-Bold.ttf'),
    ('Inter Light',    'Inter-Light.ttf',          None),
    ('Inter SemiBold', 'Inter-SemiBold.ttf',       None),
    ('JetBrains Mono', 'JetBrainsMono-Regular.ttf', None),
]

def ensure_embed_fonts():
    if all(os.path.exists(os.path.join(EMBED, f)) for _, r, b in EMBED_FILES for f in [r] + ([b] if b else [])):
        return
    from fontTools.ttLib import TTFont
    from fontTools.varLib.instancer import instantiateVariableFont
    os.makedirs(EMBED, exist_ok=True)
    src_fonts = os.path.join(HERE, '../../assets/fonts')

    def set_names(font, family, sub, bold=False):
        name = font['name']
        ps = (family + '-' + sub).replace(' ', '')
        for nid, val in [(1, family), (2, sub), (3, ps + ':BM-Embed'), (4, f'{family} {sub}'), (6, ps)]:
            name.setName(val, nid, 3, 1, 0x409)
        for nid in (16, 17):
            name.removeNames(nameID=nid)
        os2, head = font['OS/2'], font['head']
        if bold:
            os2.fsSelection = (os2.fsSelection & ~0x40) | 0x20
            head.macStyle |= 0x01
        else:
            os2.fsSelection = (os2.fsSelection | 0x40) & ~0x21
            head.macStyle &= ~0x03

    def make_inter(weight, family, sub, outname, bold=False):
        f = TTFont(os.path.join(src_fonts, 'Inter-400-latin.woff2'))
        f.flavor = None
        instantiateVariableFont(f, {'wght': weight}, inplace=True)
        f['OS/2'].usWeightClass = weight
        set_names(f, family, sub, bold)
        f.save(os.path.join(EMBED, outname))

    make_inter(400, 'Inter', 'Regular', 'Inter-Regular.ttf')
    make_inter(700, 'Inter', 'Bold', 'Inter-Bold.ttf', bold=True)
    make_inter(300, 'Inter Light', 'Regular', 'Inter-Light.ttf')
    make_inter(600, 'Inter SemiBold', 'Regular', 'Inter-SemiBold.ttf')
    jm = TTFont(os.path.join(src_fonts, 'JetBrainsMono-400-latin.woff2'))
    jm.flavor = None
    if 'fvar' in jm:
        instantiateVariableFont(jm, {'wght': 400}, inplace=True)
    set_names(jm, 'JetBrains Mono', 'Regular')
    jm.save(os.path.join(EMBED, 'JetBrainsMono-Regular.ttf'))
    print('Embed-Fonts instanziert →', EMBED)

def family_for(fw, mono):
    if mono: return 'JetBrains Mono'
    w = int(fw) if str(fw).isdigit() else (700 if fw == 'bold' else 400)
    if w <= 300: return 'Inter Light'
    if 500 <= w < 700: return 'Inter SemiBold'
    return 'Inter'

def weight_of(fw):
    return int(fw) if str(fw).isdigit() else (700 if fw == 'bold' else 400)

geo = json.load(open(os.path.join(HERE, 'geometry.json')))
wm = {v: crop_alpha(os.path.join(HERE, f'assets/wordmark-{v}.png'),
                    os.path.join(HERE, f'assets/wordmark-{v}-crop.png'))
      for v in ('black', 'white')}

prs = Presentation()
prs.slide_width, prs.slide_height = Emu(int(1440 * PX)), Emu(int(810 * PX))
blank = prs.slide_layouts[6]

ALIGN = {'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT, 'end': PP_ALIGN.RIGHT}

# Meta-Zeilen (Kopfzeile, Fußzeile, Eyebrows, Seitenzahl) bekommen
# Designbreite statt Mustertext-Maß: Die vermessene Box ist exakt so
# breit wie der Beispieltext — ein längerer echter Kapiteltitel bricht
# darin drei-, vierzeilig um (Befund Kollegen-Test 22.07.).
META_MIN_W = {'head': 560, 'eyebrow': 520, 'foot': 460, 'pagenum': 64}

def prep_content(t):
    # Voll-Blöcke (nur Inline-Kinder) tragen den kompletten Text,
    # sonst nur die eigenen Textknoten (Kinder kommen als eigene Boxen).
    content = t['text'] if t.get('full') or t['leaf'] else t['own']
    if not content: return None
    return content.upper() if t['caps'] else content


def widen_box(g, groups, s):
    """Textbox auf verfügbare Breite weiten. Die vermessene Box ist
    exakt so breit wie der Mustertext — PowerPoint bricht mit eigener
    Textmetrik minimal früher und macht aus zwei Zeilen drei (Befund
    Nils 24.07.: Cover-Titel, Kapiteltitel in der Fußzeile, unnötige
    Headline-Umbrüche). Meta-Boxen bekommen mindestens Designbreite,
    Content-Boxen die volle verfügbare Breite. Grenzen: Nachbar-Texte
    und Grafiken derselben Zeile, die eigene Karte (innen), sonst der
    Content-Rand (Meta: Seitenrand)."""
    role = g['role']
    y0, y1 = g['y'], g['y'] + g['h']
    left_lim, right_lim = 24, (1416 if role else 1320)
    for o in list(groups) + list(s.get('pix', [])):
        if o is g or not (o['y'] < y1 and o['y'] + o['h'] > y0): continue
        if o['x'] >= g['x'] + g['w']: right_lim = min(right_lim, o['x'] - 16)
        if o['x'] + o['w'] <= g['x']: left_lim = max(left_lim, o['x'] + o['w'] + 16)
    for r in s.get('rects', []):
        if (r['x'] - 8 <= g['x'] and g['x'] + g['w'] <= r['x'] + r['w'] + 8
                and r['y'] - 8 <= y0 and y1 <= r['y'] + r['h'] + 8):
            left_lim = max(left_lim, r['x'] + 20)
            right_lim = min(right_lim, r['x'] + r['w'] - 20)
    align = g['parts'][0]['align']
    if role:
        min_w = META_MIN_W[role]
        if g['w'] >= min_w: return g['x'], g['w']
        if align in ('right', 'end'):
            r = g['x'] + g['w']
            x = max(left_lim, r - min_w)
            return x, max(r - x, g['w'])
        if align == 'center':
            cx = g['x'] + g['w'] / 2
            half = min(min_w / 2, cx - left_lim, right_lim - cx)
            return cx - half, max(half * 2, g['w'])
        return g['x'], max(min(min_w, right_lim - g['x']), g['w'])
    if align in ('right', 'end'):
        r = g['x'] + g['w']
        x = min(g['x'], left_lim)
        return x, r - x
    if align == 'center':
        cx = g['x'] + g['w'] / 2
        half = min(cx - left_lim, right_lim - cx)
        return (cx - half, half * 2) if half * 2 > g['w'] else (g['x'], g['w'])
    return g['x'], max(g['w'], right_lim - g['x'])

def merge_columns(texts):
    """Gestapelte Content-Boxen einer Spalte (Titel/Body/Caption) werden
    EIN Textrahmen: PPTX-Boxen schieben sich nicht weg wie HTML-Blöcke —
    läuft echter Text länger als das Muster, wächst er sonst aus seiner
    Box direkt in die nächste (Text-über-Text im Kollegen-Test)."""
    items = []
    for t in texts:
        c = prep_content(t)
        if c is None: continue
        items.append({'x': t['x'], 'y': t['y'], 'w': t['w'], 'h': t['h'],
                      'role': t.get('role'), 'parts': [{**t, 'content': c, 'gap': 0}]})
    items.sort(key=lambda g: (round(g['x']), g['y']))
    merged = []
    for g in items:
        prev = merged[-1] if merged else None
        last = prev['parts'][-1] if prev else None
        gap = g['y'] - (prev['y'] + prev['h']) if prev else None
        if (prev and not g['role'] and not prev['role']
                and g['parts'][0]['fs'] <= 20 and last['fs'] <= 20
                and abs(g['x'] - prev['x']) < 4
                and -2 <= gap <= 32
                and g['parts'][0]['align'] == last['align']):
            g['parts'][0]['gap'] = max(gap, 0)
            prev['parts'] += g['parts']
            prev['w'] = max(prev['w'], g['w'])
            prev['h'] = g['y'] + g['h'] - prev['y']
        else:
            merged.append(g)
    return merged

ensure_embed_fonts()

for i, s in enumerate(geo['slides']):
    slide = prs.slides.add_slide(blank)

    # Grund: Moment-Slides tragen das gerenderte Linien-/Kasane-Bild,
    # Content-Slides die Off-White-Fläche. Die Hintergründe gehen als
    # JPEG in die Datei: Der Grain-Layer ist Rauschen — als PNG wird
    # jede Fläche 2-3 MB (14-MB-Template), als JPEG ~300 KB, ohne
    # sichtbaren Unterschied auf der Vollfläche (kein Alpha nötig).
    if s['moment']:
        bg_png = os.path.join(HERE, f'assets/bg-{i:02d}.png')
        bg_jpg = os.path.join(HERE, f'assets/bg-{i:02d}.jpg')
        Image.open(bg_png).convert('RGB').save(bg_jpg, quality=90)
        slide.shapes.add_picture(bg_jpg, 0, 0, prs.slide_width, prs.slide_height)
    else:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(s['bg'])

    # Flächen (Karten, Bänder, Rahmen) und Grafik-Ausschnitte (SVGs,
    # Gradients) in DOM-Reihenfolge — verschachtelte Elemente liegen
    # so korrekt übereinander.
    layers = sorted(
        [{'t': 'rect', **r} for r in s.get('rects', [])]
        + [{'t': 'pix', **p} for p in s.get('pix', [])],
        key=lambda o: o.get('z', 0))
    shot_img = None
    for o in layers:
        if o['w'] < 2 or o['h'] < 2: continue
        if o['t'] == 'rect':
            kind = MSO_SHAPE.ROUNDED_RECTANGLE if o.get('radius', 0) > 2 else MSO_SHAPE.RECTANGLE
            shp = flat(slide.shapes.add_shape(kind, emu(o['x']), emu(o['y']), emu(o['w']), emu(o['h'])))
            if kind == MSO_SHAPE.ROUNDED_RECTANGLE:
                shp.adjustments[0] = max(0.0, min(0.5, o['radius'] / min(o['w'], o['h'])))
            shp.fill.fore_color.rgb = rgb(o['color'])
        else:
            if shot_img is None:
                shot_img = Image.open(os.path.join(SHOTS, f'slide-{i:02d}.png'))
            box = (int(o['x'] * 2), int(o['y'] * 2), int((o['x'] + o['w']) * 2), int((o['y'] + o['h']) * 2))
            cp = os.path.join(SHOTS, f'pix-{i:02d}-{o.get("z", 0)}.png')
            shot_img.crop(box).save(cp)
            slide.shapes.add_picture(cp, emu(o['x']), emu(o['y']), emu(o['w']), emu(o['h']))

    # Hairlines als flache Shapes (dedupliziert)
    seen = set()
    for ln in s['lines']:
        key = (round(ln['x']), round(ln['y']), round(ln['w']), round(ln['h']))
        if key in seen: continue
        seen.add(key)
        shp = flat(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(ln['x']), emu(ln['y']),
                                          emu(max(ln['w'], 1)), emu(max(ln['h'], 1))))
        shp.fill.fore_color.rgb = rgb(ln['color'])

    # Wortmarken
    for img in s['images']:
        v = img['kind'].split('-')[1]
        slide.shapes.add_picture(wm[v], emu(img['x']), emu(img['y']),
                                 emu(img['w']), emu(img['h']))

    # Texte: eigene Textknoten am eigenen Element (Duplikate über
    # Eltern/Kind-Verschachtelung vermeidet der own-Ansatz);
    # Spalten-Stapel vorab zu einem Rahmen zusammengelegt.
    groups = merge_columns(s['texts'])
    for grp in groups:
        role = grp['role']
        x, y, w, h = grp['x'], grp['y'], grp['w'], grp['h']
        x, w = widen_box(grp, groups, s)
        pad_y = 4 if role else 2
        tb = slide.shapes.add_textbox(emu(x), emu(y - pad_y), emu(w + 4), emu(h + 2 * pad_y))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        if role:
            tf.word_wrap = False                          # Meta bleibt einzeilig
            tf.auto_size = MSO_AUTO_SIZE.NONE
            # Mittig ankern: PowerPoints Zeilenbox weicht von Chromes
            # engem Text-Rechteck ab — top-verankert hing „[ Kundenlogo ]"
            # sichtbar über der Brückenlinie (Befund Nils 24.07.).
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        else:
            tf.word_wrap = True
            # normAutofit: läuft echter Text länger als das Muster,
            # schrumpft er in der Box, statt darüber hinauszuwachsen.
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        li = 0
        for part in grp['parts']:
            for pi, line in enumerate(part['content'].split('\n')):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                li += 1
                p.text = line
                if part['align'] in ALIGN: p.alignment = ALIGN[part['align']]
                # Zeilenabstand = vermessene HTML-Zeilenhöhe; ohne sie
                # rendert PowerPoint seinen Default und der Text sitzt
                # in jeder mehrzeiligen Box falsch (Befund 23.07.).
                if part.get('lh'):
                    p.line_spacing = Pt(round(part['lh'] * 2 / 3, 1))
                if pi == 0 and part['gap']:
                    p.space_before = Pt(round(part['gap'] * 2 / 3, 1))
                f = p.runs[0].font
                f.name = family_for(part['fw'], part['mono'])
                f.size = Pt(round(part['fs'] * 2 / 3, 1))
                f.bold = weight_of(part['fw']) >= 700
                f.italic = part['italic']
                f.color.rgb = rgb(part['color'])
                # Laufweite (negatives Headline-Tracking, Eyebrow-
                # Sperrung) als a:spc in 1/100 pt.
                mls = re.match(r'(-?[\d.]+)px', part.get('ls') or '')
                if mls:
                    spc = int(round(float(mls.group(1)) * 2 / 3 * 100))
                    if spc:
                        p.runs[0]._r.get_or_add_rPr().set('spc', str(spc))

    # Sprechernotiz = Layout-Name (hilft Langdocks Klassifizierung)
    slide.notes_slide.notes_text_frame.text = f"Layout: {s['label']}"

prs.save(OUT)

# ---------- Schriften einbetten (Zip-Nachbearbeitung) ----------
# python-pptx kann keine Fonts einbetten; die Teile werden direkt
# ins OOXML-Paket geschrieben: fntdata-Parts, Content-Type,
# Relationships, p:embeddedFontLst + embedTrueTypeFonts.
FONT_REL = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/font'
zin = zipfile.ZipFile(OUT)
items = {n: zin.read(n) for n in zin.namelist()}
zin.close()

ct = items['[Content_Types].xml'].decode('utf-8')
if 'fntdata' not in ct:
    ct = ct.replace('</Types>', '<Default Extension="fntdata" ContentType="application/x-fontdata"/></Types>')
items['[Content_Types].xml'] = ct.encode('utf-8')

rel_add, embeds, fi = '', [], 1
for fam, reg, bold in EMBED_FILES:
    entry = f'<p:embeddedFont><p:font typeface="{fam}"/>'
    for slot, fn in (('regular', reg), ('bold', bold)):
        if not fn: continue
        rid = f'rIdFont{fi}'
        items[f'ppt/fonts/font{fi}.fntdata'] = open(os.path.join(EMBED, fn), 'rb').read()
        rel_add += f'<Relationship Id="{rid}" Type="{FONT_REL}" Target="fonts/font{fi}.fntdata"/>'
        entry += f'<p:{slot} r:id="{rid}"/>'
        fi += 1
    embeds.append(entry + '</p:embeddedFont>')

rels = items['ppt/_rels/presentation.xml.rels'].decode('utf-8')
items['ppt/_rels/presentation.xml.rels'] = rels.replace('</Relationships>', rel_add + '</Relationships>').encode('utf-8')

pres = items['ppt/presentation.xml'].decode('utf-8')
if 'embedTrueTypeFonts' not in pres:
    pres = pres.replace('<p:presentation ', '<p:presentation embedTrueTypeFonts="1" ', 1)
lst = '<p:embeddedFontLst>' + ''.join(embeds) + '</p:embeddedFontLst>'
pres, n = re.subn(r'(<p:notesSz[^>]*/>)', r'\1' + lst, pres, count=1)
if not n:
    raise SystemExit('FEHLER — notesSz nicht gefunden, Font-Liste nicht eingefügt.')
items['ppt/presentation.xml'] = pres.encode('utf-8')

with zipfile.ZipFile(OUT, 'w', zipfile.ZIP_DEFLATED) as z:
    for n, d in items.items():
        z.writestr(n, d)

print(f'OK — {len(geo["slides"])} Folien, {os.path.getsize(OUT) // 1024} KB (Fonts eingebettet) -> {OUT}')
